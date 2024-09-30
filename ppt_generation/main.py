from pptx import Presentation
from pptx.util import Inches
import streamlit as st

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # Layout for the title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Your Presentation Title"
    subtitle.text = "Subtitle or Description"

    return prs

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Layout for a content slide
    slide = prs.slides.add_slide(slide_layout)
    slide_title = slide.shapes.title
    text_box = slide.shapes.placeholders[1]

    slide_title.text = title  # Set the title
    text_box.text = content  # Set the content as bullet points

def text_to_slides(prs, text):
    # Split the input text into sections
    sections = text.split('\n\n')  # Split by paragraphs (or any delimiter)
    
    for i, section in enumerate(sections):
        # Assume the first line is the title, the rest is content
        lines = section.split('\n')
        title = lines[0] if len(lines) > 0 else f"Slide {i+1}"
        content = '\n'.join(lines[1:]) if len(lines) > 1 else ""
        
        add_slide(prs, title, content)

def save_presentation(prs):
    # Save the presentation to a BytesIO object
    from io import BytesIO
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def add_image_slide(prs, image_path, slide_title):
    slide_layout = prs.slide_layouts[5]  # Title and image layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_title

    img = slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(4))

# Streamlit UI
st.title("Text-to-PPT Generator")

text_input = st.text_area("Enter your text here:", height=300)

if st.button("Generate PPT"):
    prs = create_presentation()
    text_to_slides(prs, text_input)
    ppt_buffer = save_presentation(prs)  # Save to a BytesIO object
    st.success("Presentation generated!")
    
    # Use st.download_button for downloading
    st.download_button(
        label="Download PPT",
        data=ppt_buffer,
        file_name="output.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
